#!/usr/bin/env python3
"""
LGOX多格式文档摄入管道 v1.0
端口:8770 层:L0→L1(L1知识层)
支持: PDF/Word/PPT/Markdown/TXT/HTML → 解析→清洗→LGE基因
基因ID: GENE-PRO-doc-ingest-v1
"""
import json, os, sys, io, hashlib, time, traceback
from http.server import HTTPServer, BaseHTTPRequestHandler
from urllib.parse import parse_qs, urlparse
from pathlib import Path

LGE_URL = "http://100.116.0.29:8200"
DOCS_DIR = Path(os.path.expanduser("~/lgox-docs"))
DOCS_DIR.mkdir(parents=True, exist_ok=True)

def parse_pdf(filepath: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    except Exception as e:
        return f"[PDF解析失败: {e}]"

def parse_docx(filepath: str) -> str:
    try:
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        return f"[DOCX解析失败: {e}]"

def parse_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n---\n".join(texts)
    except Exception as e:
        return f"[PPTX解析失败: {e}]"

def parse_markdown(filepath: str) -> str:
    with open(filepath) as f:
        return f.read()

def parse_txt(filepath: str) -> str:
    with open(filepath, encoding='utf-8', errors='replace') as f:
        return f.read()

PARSERS = {
    '.pdf': parse_pdf, '.docx': parse_docx, '.pptx': parse_pptx,
    '.md': parse_markdown, '.txt': parse_txt, '.html': parse_txt,
    '.py': parse_txt, '.json': parse_txt, '.yaml': parse_txt,
    '.yml': parse_txt, '.csv': parse_txt
}

def write_gene(content: str, source_file: str, source_type: str) -> dict:
    """写入LGE基因引擎"""
    import urllib.request
    gene_id = f"GENE-DOC-{hashlib.sha256(content.encode()).hexdigest()[:12]}"
    payload = json.dumps({
        "gene_id": gene_id,
        "content": content[:8000],
        "source_file": source_file,
        "source_type": source_type,
        "content_length": len(content),
        "category": "document",
        "domain": "general",
        "tags": ["doc-ingest", source_type],
        "quality_score": 0.7
    }).encode()
    try:
        req = urllib.request.Request(f"{LGE_URL}/genes/write", data=payload,
            headers={"Content-Type": "application/json"}, method="POST")
        with urllib.request.urlopen(req, timeout=10) as resp:
            result = json.loads(resp.read())
        return {"gene_id": gene_id, "lge_status": result.get("status", "ok"), "content_len": len(content)}
    except Exception as e:
        return {"gene_id": gene_id, "lge_status": f"write_failed: {e}", "content_len": len(content)}

class DocIngestHandler(BaseHTTPRequestHandler):
    def do_GET(self):
        path = urlparse(self.path).path
        if path == "/health":
            self._json({"status":"ok","service":"LGOX文档摄入管道","version":"1.0","layer":"L1知识层(Know)"})
        elif path == "/stats":
            files = list(DOCS_DIR.glob("*"))
            self._json({"total_docs": len(files), "docs_dir": str(DOCS_DIR)})
        else:
            self._json({"error":"use POST /ingest"})

    def do_POST(self):
        path = urlparse(self.path).path
        if path == "/ingest":
            content_len = int(self.headers.get('Content-Length', 0))
            if content_len == 0:
                self._json({"error": "no content"}, 400)
                return

            body = self.rfile.read(content_len)
            content_type = self.headers.get('Content-Type', '')

            if 'multipart/form-data' in content_type:
                self._json({"error": "multipart not supported, use raw POST"}, 400)
                return

            # Try JSON with base64 file
            try:
                data = json.loads(body)
                filename = data.get('filename', 'uploaded.bin')
                content = data.get('content', '')
                if not content:
                    self._json({"error": "no content"}, 400)
                    return

                # Save raw
                ext = Path(filename).suffix.lower()
                saved_path = DOCS_DIR / f"{int(time.time())}_{filename}"
                saved_path.write_text(content, encoding='utf-8', errors='replace')

                # Parse
                parser = PARSERS.get(ext, parse_txt)
                text = parser(str(saved_path))

                # Write LGE gene
                result = write_gene(text, filename, ext.lstrip('.'))

                self._json({
                    "file": filename,
                    "saved": str(saved_path),
                    "text_len": len(text),
                    "gene_id": result.get("gene_id"),
                    "lge": result.get("lge_status")
                })
            except json.JSONDecodeError:
                # Raw text upload
                saved_path = DOCS_DIR / f"{int(time.time())}_raw.txt"
                saved_path.write_text(body.decode('utf-8', errors='replace'))
                result = write_gene(body.decode('utf-8', errors='replace'), 'raw.txt', 'txt')
                self._json({
                    "file": "raw.txt", "saved": str(saved_path),
                    "text_len": len(body),
                    "gene_id": result.get("gene_id"),
                    "lge": result.get("lge_status")
                })
        else:
            self._json({"error": "not found"}, 404)

    def _json(self, data, code=200):
        self.send_response(code)
        self.send_header('Content-Type', 'application/json')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode())

    def log_message(self, format, *args):
        pass

if __name__ == '__main__':
    port = int(sys.argv[1]) if len(sys.argv) > 1 else 8770
    server = HTTPServer(('127.0.0.1', port), DocIngestHandler)
    print(f"📄 文档摄入管道 :{port} → LGE(地枢:8200)")
    server.serve_forever()
